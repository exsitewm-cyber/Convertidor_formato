import io
import streamlit as st
from pptx import Presentation
from pptx.util import Inches
from lxml import etree

st.set_page_config(page_title="Editor PPTX Avanzado", page_icon="📊", layout="wide")

# Inicializar memoria
if 'plantillas_guardadas' not in st.session_state:
    st.session_state['plantillas_guardadas'] = {}

# ==========================================
# 1. EXTRACTOR DE TEMA Y ESTILOS (Para el Preview)
# ==========================================
def extraer_tema_plantilla(bytes_plantilla):
    """Lee el XML interno de la plantilla para extraer fuentes y paleta de colores reales"""
    theme_data = {
        'font_title': 'sans-serif',
        'font_body': 'sans-serif',
        'color_title': '#000000',
        'color_body': '#333333',
        'bg_color': '#ffffff'
    }
    try:
        prs = Presentation(io.BytesIO(bytes_plantilla))
        for part in prs.part.package.parts:
            # Buscar la parte del Tema (Theme)
            if part.partname.startswith('/ppt/theme/'):
                root = etree.fromstring(part.blob)
                ns = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
                
                # Extraer Fuentes Maestras
                maj_font = root.xpath('//a:majorFont/a:latin/@typeface', namespaces=ns)
                if maj_font: theme_data['font_title'] = maj_font[0]
                min_font = root.xpath('//a:minorFont/a:latin/@typeface', namespaces=ns)
                if min_font: theme_data['font_body'] = min_font[0]
                
                # Extraer Colores Maestros (Oscuro 1 para textos)
                dk1 = root.xpath('//a:clrScheme/a:dk1/a:srgbClr/@val', namespaces=ns)
                if dk1: 
                    theme_data['color_title'] = '#' + dk1[0]
                    theme_data['color_body'] = '#' + dk1[0]
                else:
                    sys_dk1 = root.xpath('//a:clrScheme/a:dk1/a:sysClr/@lastClr', namespaces=ns)
                    if sys_dk1:
                        theme_data['color_title'] = '#' + sys_dk1[0]
                        theme_data['color_body'] = '#' + sys_dk1[0]
                        
                # Color de fondo (Claro 1)
                lt1 = root.xpath('//a:clrScheme/a:lt1/a:srgbClr/@val', namespaces=ns)
                if lt1: theme_data['bg_color'] = '#' + lt1[0]
                break
    except Exception as e:
        pass
    return theme_data

# ==========================================
# 2. EXTRACTOR DE CONTENIDO ORIGINAL
# ==========================================
def extraer_secciones_pptx(file_stream, max_slides=5):
    """Extrae el texto puro dividiendo rigurosamente Título vs Cuerpo"""
    file_stream.seek(0)
    prs = Presentation(file_stream)
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        if i >= max_slides: break
        slide_info = {"titulo": "", "cuerpo": []}
        
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text.strip():
                continue
            
            # Si es el marcador de Título
            if shape == slide.shapes.title:
                slide_info["titulo"] = shape.text
            else:
                # Si es contenido, guardamos el texto y su nivel de viñeta
                for p in shape.text_frame.paragraphs:
                    if p.text.strip():
                        slide_info["cuerpo"].append({'text': p.text, 'level': p.level})
                        
        slides_data.append(slide_info)
    file_stream.seek(0)
    return slides_data

# ==========================================
# 3. PROCESADOR Y APLICADOR DE FORMATO
# ==========================================
def procesar_pptx(stream_origen, stream_plantilla_bytes):
    """Inyecta el texto limpio en la plantilla forzando la adopción del formato maestro"""
    stream_origen.seek(0)
    prs_origen = Presentation(stream_origen)
    prs_plantilla = Presentation(io.BytesIO(stream_plantilla_bytes))

    # Limpiar diapositivas de muestra de la plantilla
    xml_slides = prs_plantilla.slides._sldIdLst  
    for s in list(xml_slides):
        xml_slides.remove(s)

    for slide_orig in prs_origen.slides:
        # Detectar el diseño original o usar Título y Objetos por defecto
        layout_idx = 1 
        if slide_orig.slide_layout.name == "Title Slide": layout_idx = 0
        if layout_idx >= len(prs_plantilla.slide_layouts): layout_idx = 1
        
        layout_target = prs_plantilla.slide_layouts[layout_idx]
        new_slide = prs_plantilla.slides.add_slide(layout_target)
        
        # Extraer datos de la diapositiva original
        titulo = ""
        cuerpo = []
        for shape in slide_orig.shapes:
            if not shape.has_text_frame or not shape.text.strip(): continue
            if shape == slide_orig.shapes.title:
                titulo = shape.text
            else:
                for p in shape.text_frame.paragraphs:
                    if p.text.strip(): cuerpo.append({'text': p.text, 'level': p.level})

        # INYECCIÓN PURA: Dejar que PowerPoint aplique su formato
        # 1. Título
        if new_slide.shapes.title and titulo:
            new_slide.shapes.title.text = titulo  # Asignación directa hereda el formato
            
        # 2. Cuerpo
        if cuerpo:
            body_ph = None
            for ph in new_slide.placeholders:
                if ph.placeholder_format.idx != 0: # Buscar un marcador que no sea título
                    body_ph = ph
                    break
            
            if body_ph:
                tf = body_ph.text_frame
                tf.clear() # Borramos para resetear el formato base del marcador
                for i, p_data in enumerate(cuerpo):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = p_data['text']
                    try: p.level = p_data['level'] # Mantener nivel de viñeta
                    except: pass
            else:
                # Si la plantilla no tiene cuadro de cuerpo, creamos uno básico
                tb = new_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
                tf = tb.text_frame
                for i, p_data in enumerate(cuerpo):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = p_data['text']

    output_stream = io.BytesIO()
    prs_plantilla.save(output_stream)
    output_stream.seek(0)
    return output_stream

# ==========================================
# INTERFAZ (UI) - GESTOR DE PLANTILLAS
# ==========================================
st.sidebar.header("📁 Gestor de Plantillas")
if st.session_state['plantillas_guardadas']:
    for nombre_plantilla in list(st.session_state['plantillas_guardadas'].keys()):
        col1, col2 = st.sidebar.columns([0.8, 0.2])
        col1.write(f"🎨 {nombre_plantilla}")
        if col2.button("❌", key=f"del_{nombre_plantilla}"):
            del st.session_state['plantillas_guardadas'][nombre_plantilla]
            st.rerun()

st.sidebar.markdown("---")
if len(st.session_state['plantillas_guardadas']) < 5:
    nuevas_plantillas = st.sidebar.file_uploader("Añadir plantilla base (.pptx)", type=["pptx"], accept_multiple_files=True)
    if nuevas_plantillas:
        for p in nuevas_plantillas:
            if len(st.session_state['plantillas_guardadas']) < 5:
                st.session_state['plantillas_guardadas'][p.name] = p.getvalue()
        st.rerun()

# ==========================================
# ÁREA PRINCIPAL Y PREVISUALIZACIÓN DINÁMICA
# ==========================================
st.title("📊 Editor Automático y Previsualizador de PPTX")

archivo_original = st.file_uploader("1. Sube tu archivo a modificar (.pptx)", type=["pptx"])

if archivo_original and st.session_state['plantillas_guardadas']:
    st.markdown("---")
    
    plantilla_seleccionada = st.selectbox("2. Elige la Plantilla a aplicar:", list(st.session_state['plantillas_guardadas'].keys()))
    bytes_plantilla = st.session_state['plantillas_guardadas'][plantilla_seleccionada]
    
    # Extraer colores y fuentes de la plantilla seleccionada para inyectarlos en el HTML
    tema_plantilla = extraer_tema_plantilla(bytes_plantilla)
    datos_diapositivas = extraer_secciones_pptx(archivo_original)
    
    st.markdown(f"### 👀 Vista Previa con el Diseño: *{plantilla_seleccionada}*")
    st.caption("Esta vista lee la tipografía y los colores internos de tu plantilla y te muestra cómo se distribuirán tus secciones.")
    
    # CSS Dinámico basado en la plantilla seleccionada
    bg_color = tema_plantilla['bg_color'] if tema_plantilla['bg_color'] != '#000000' else '#ffffff'
    
    for idx, slide in enumerate(datos_diapositivas):
        st.markdown(f"**Diapositiva {idx+1}**")
        
        # Generar HTML simulando la plantilla real
        html_preview = f"""
        <div style="background-color: {bg_color}; border: 1px solid #ddd; padding: 25px; border-radius: 10px; margin-bottom: 20px; box-shadow: 2px 2px 10px rgba(0,0,0,0.1);">
            <div style="margin-bottom: 20px;">
                <span style="background-color: #444; color: white; padding: 3px 8px; border-radius: 4px; font-size: 11px; margin-bottom: 5px; display: inline-block;">MARCADOR: TÍTULO</span>
                <h2 style="font-family: '{tema_plantilla['font_title']}', sans-serif; color: {tema_plantilla['color_title']}; margin: 0; padding-bottom: 10px; border-bottom: 2px solid {tema_plantilla['color_title']}33;">
                    {slide['titulo'] if slide['titulo'] else '<i>[Sin Título]</i>'}
                </h2>
            </div>
            <div>
                <span style="background-color: #444; color: white; padding: 3px 8px; border-radius: 4px; font-size: 11px; margin-bottom: 5px; display: inline-block;">MARCADOR: CUERPO Y VIÑETAS</span>
                <div style="font-family: '{tema_plantilla['font_body']}', sans-serif; color: {tema_plantilla['color_body']}; font-size: 16px; margin-top: 10px;">
        """
        
        # Procesar los niveles de viñetas para la vista previa web
        if slide['cuerpo']:
            for item in slide['cuerpo']:
                indent = item['level'] * 20 if item['level'] else 0
                html_preview += f"<div style='margin-left: {indent}px;'>• {item['text']}</div>"
        else:
            html_preview += "<i>[Contenido vacío]</i>"
            
        html_preview += "</div></div></div>"
        
        st.markdown(html_preview, unsafe_allow_html=True)

    st.markdown("---")
    
    # Procesamiento y Descarga
    if st.button("🚀 Aplicar Plantilla y Descargar", type="primary"):
        with st.spinner(f"Aplicando formato estricto de '{plantilla_seleccionada}'..."):
            ppt_final = procesar_pptx(archivo_original, bytes_plantilla)
            
        st.success("¡Documento generado aplicando la fuente, colores, tamaño y distribución de la plantilla!")
        st.download_button(
            label="📥 Descargar Documento Formateado",
            data=ppt_final,
            file_name=f"Plantilla_{plantilla_seleccionada}_{archivo_original.name}",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
